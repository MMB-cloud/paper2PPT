import os

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from common.Utils import Utils
from pptx.util import Inches

utils = Utils()

"""实现对幻灯片模板内容的读取与写入"""


def replace_text(text_frame, title):
    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(title)
    pass


def set_style(text_frame, style):
    paragraph = text_frame.paragraphs[0]
    font = paragraph.font
    font.name = style["font_name"]
    font.color.rgb = RGBColor(style["font_color_rgb"][0], style["font_color_rgb"][1], style["font_color_rgb"][2])
    font.bold = style["font_bold"]
    font.size = Pt(style["font_size"])
    if style["paragraph_alignment"] == "center":
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    pass


def add_text_style(slide, content, curSentCnt, style):
    # 设置文本框的宽度和高度
    text_box_width = Inches(10)
    text_box_height = Inches(1.5)
    # 计算文本框的起始位置
    left = Inches(1.5)
    top = Inches(0.5 + (curSentCnt * 1.5 + 0.2))  # + (curSentCnt * text_box_height + 0.1)

    textbox = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
    # 获取文本框内的文本框架
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.autofit = True
    # 设置段落的内容
    tf.text = content
    set_style(tf, style)


def add_image(image_path, slide):
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(4.5)
    slide.shapes.add_picture(image_path, left, top, width, height)


def add_tables(slide, table_data, rows, cols):
    # 添加表格
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(4.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置单元格内容
    for row in range(len(table_data)):
        for col in range(len(table_data[row])):  # 存在表格中列数不一致
            cell = table.cell(row, col)
            cell.text = table_data[row][col]


class Template:
    template_dir_path = ''
    output_path = ''
    configure = {}
    """
    :param template_dir_path: 幻灯片模板文件夹
    :param output_dir_path: 幻灯片输出文件夹
    :return: None
    """

    def __init__(self, template_dir_path, output_path) -> None:
        self.template_dir_path = template_dir_path
        self.output_path = output_path
        self.configure = utils.json_to_dict(self.template_dir_path + "\\template_01" + "\\configure.json")

    def findPathByRid(self, rid):
        for file in os.listdir(self.output_path):
            if str(rid) in file:
                return file
        return -1

    """
    :param data: 待输出的文本、图表和图像内容
    :param output_path: 替换后幻灯片的指定输出路径
    """

    def replace(self, data, output_path):
        # 1. 根据slide类型确定要替换的pptx_path
        slide_index = -1
        if data["type"] == 0:
            slide_index = 0
        elif data["type"] == 1:
            num = len(data["content"])
            slide_index = num - 3
        elif data["type"] == 2:
            slide_index = 4
        elif data["type"] == 3:
            slide_index = 5
        pptx_path = self.template_dir_path + "\index_0" + str(slide_index) + ".pptx"
        ppt = Presentation(pptx_path)
        slide = ppt.slides[0]

        # 2. 获取pptx的配置
        configure = utils.json_to_dict(self.template_dir_path + "\configure.json")
        text_frames_configure = configure["slides"][slide_index]["text_frames"]

        # 3. 获取文本框，并进行内容替换
        text_frames = [] if slide_index not in [5] else {}
        if slide_index not in [5]:
            for text_frame_configure in text_frames_configure:
                for shape in slide.shapes:
                    if shape.has_text_frame and text_frame_configure["content"] in shape.text_frame.text:
                        text_frames.append(shape.text_frame)
            for i in range(len(text_frames)):
                text_frame = text_frames[i]
                text_frame_style = text_frames_configure[i]["style"]
                paragraph = text_frame.paragraphs[0]
                paragraph.text = data["content"][i] if type(data["content"]) == list else data["content"]
                if text_frame_style["paragraph_alignment"] == "center":
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                font = paragraph.font
                font.name = text_frame_style["font_name"]
                font.color.rgb = RGBColor(text_frame_style["font_color_rgb"][0], text_frame_style["font_color_rgb"][1],
                                          text_frame_style["font_color_rgb"][2])
                font.size = Pt(text_frame_style["font_size"])
                font.bold = text_frame_style["font_bold"]
        else:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if text_frames_configure["heading"]["content"] in shape.text_frame.text:
                        paragraph = shape.text_frame.paragraphs[0]
                        text_frame_style = text_frames_configure["heading"]["style"]
                        paragraph.text = data["content"]["heading"]
                        if text_frame_style["paragraph_alignment"] == "center":
                            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        font = paragraph.font
                        font.name = text_frame_style["font_name"]
                        font.color.rgb = RGBColor(text_frame_style["font_color_rgb"][0],
                                                  text_frame_style["font_color_rgb"][1],
                                                  text_frame_style["font_color_rgb"][2])
                        font.size = Pt(text_frame_style["font_size"])
                        font.bold = text_frame_style["font_bold"]
                    else:
                        for i in range(len(text_frames_configure["children"])):
                            if text_frames_configure["children"][i]["keyword"]["content"] in shape.text_frame.text:
                                paragraph = shape.text_frame.paragraphs[0]
                                text_frame_style = text_frames_configure["children"][i]["keyword"]["style"]
                                paragraph.text = data["content"]["children"][i]["keyword"]
                                if text_frame_style["paragraph_alignment"] == "center":
                                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                                font = paragraph.font
                                font.name = text_frame_style["font_name"]
                                font.color.rgb = RGBColor(text_frame_style["font_color_rgb"][0],
                                                          text_frame_style["font_color_rgb"][1],
                                                          text_frame_style["font_color_rgb"][2])
                                font.size = Pt(text_frame_style["font_size"])
                                font.bold = text_frame_style["font_bold"]
                            else:
                                if text_frames_configure["children"][i]["sentences"][
                                    "content"] in shape.text_frame.text:
                                    for j in range(len(data["content"]["children"][i]["sentences"])):
                                        paragraph = shape.text_frame.paragraphs[j]
                                        text_frame_style = text_frames_configure["children"][i]["sentences"]["style"]
                                        paragraph.text = data["content"]["children"][i]["sentences"][j]
                                        if text_frame_style["paragraph_alignment"] == "center":
                                            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                                        font = paragraph.font
                                        font.name = text_frame_style["font_name"]
                                        font.color.rgb = RGBColor(text_frame_style["font_color_rgb"][0],
                                                                  text_frame_style["font_color_rgb"][1],
                                                                  text_frame_style["font_color_rgb"][2])
                                        font.size = Pt(text_frame_style["font_size"])
                                        font.bold = text_frame_style["font_bold"]
        ppt.save(output_path)

    def setTitle(self, title):
        pptx_path = self.template_dir_path + "\\template_01" + "\\index_00.pptx"
        self.output_path = self.output_path + "\\" + title
        ppt = Presentation(pptx_path)
        slide = ppt.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.paragraphs[0].text == \
                    self.configure["slides"][0]["text_frames"][0]["content"]:
                replace_text(shape.text_frame, title)
                set_style(shape.text_frame, self.configure["slides"][0]["text_frames"][0]["style"])
                break
        ppt.save(self.output_path + "\\output_0_0.pptx")

    def setIndex(self, index):
        pptx_path = self.template_dir_path + "\\template_01" + "\\index_01.pptx"
        ppt = Presentation(pptx_path)
        slide = ppt.slides[0]
        # 设置文本框的宽度和高度
        text_box_width = Inches(2.5)
        text_box_height = Inches(0.5)
        # 计算文本框的起始位置
        left = Inches(5.6)
        top = Inches(1.5)
        # 逐项添加列表中的内容到单独的文本框中
        i = 1
        for item in index:
            # 创建一个文本框
            textbox = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
            # 获取文本框内的文本框架
            tf = textbox.text_frame
            # 设置段落的内容
            tf.text = str(i) + " " + item
            i += 1
            set_style(tf, self.configure["slides"][1]["text_frames"][0]["style"])
            # 更新文本框的顶部位置
            top += text_box_height + Inches(0.1)  # 增加间距
        ppt.save(self.output_path + "\\output_0_1.pptx")

    # TODO 该以什么粒度入参？ 考虑一个part作为入参，一个part可能有多个slide
    def setSlideNode(self, pptNode, index_script, index_title):
        # 标题页 替换序号和文字内容即可
        style = {
            "paragraph_alignment": "center",
            "font_name": "微软雅黑",
            "font_color_rgb": [255, 255, 255],
            "font_size": 44,
            "font_bold": True
        }
        pptx_path = self.template_dir_path + "\\template_01" + "\\index_04.pptx"
        ppt = Presentation(pptx_path)
        slide = ppt.slides[0]
        cnt = -1
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.paragraphs[0].text.isdigit():
                for i in range(len(index_script)):
                    if pptNode.getContent() == index_script[i]:
                        cnt = i + 1
                replace_text(shape.text_frame, cnt)
                set_style(shape.text_frame, style)
            elif shape.has_text_frame:
                replace_text(shape.text_frame, index_title[cnt - 1])
                set_style(shape.text_frame, style)
        ppt.save(self.output_path + "\\output_" + str(cnt) + "_" + "0.pptx")

        # 节点页
        normal_path = self.template_dir_path + "\\template_01" + "\\index_normal.pptx"
        picture_path = self.template_dir_path + "\\template_01" + "\\index_picture.pptx"
        # 单页最大句子数量
        sentence_max_cnt = 4
        '''
        图表和公式节点每页一个
        句子节点每页四个
        计算出每个part需要的ppt数量，一次创建完，然后一页一页填充
        '''
        # normalSlides, picSlides = pptNode.getSlideCountByIndex(cnt - 1)
        normal_style = {
            "paragraph_alignment": "left",
            "font_name": "微软雅黑",
            "font_color_rgb": [0, 0, 0],
            "font_size": 20,
            "font_bold": False
        }
        title_style = {
            "paragraph_alignment": "left",
            "font_name": "微软雅黑",
            "font_color_rgb": [0, 0, 0],
            "font_size": 22,
            "font_bold": True
        }
        curSlide = 1
        curSentCnt = 0
        ppt = Presentation(normal_path)
        slide = ppt.slides[0]
        for node in pptNode.getChildren():
            # 文本节点
            if node["type"] == 2:
                slide_parts = node["content"]
                slide_sents = slide_parts["content"]
                add_text_style(slide, slide_parts['title'], curSentCnt, title_style)
                curSentCnt += 1
                for slide_sent in slide_sents:
                    add_text_style(slide, slide_sent, curSentCnt, normal_style)
                    curSentCnt += 1
                    # 够一张幻灯片了 保存 重置参数
                    if curSentCnt >= sentence_max_cnt:
                        ppt.save(self.output_path + "\\output_" + str(cnt) + "_" + str(curSlide) + ".pptx")
                        ppt = Presentation(normal_path)
                        slide = ppt.slides[0]
                        curSlide += 1
                        curSentCnt = 0
            # 图片节点
            if node["type"] == 4:
                ppt_pic = Presentation(picture_path)
                slide_pic = ppt_pic.slides[0]
                pic_path = self.findPathByRid(node["content"])
                if pic_path == -1:
                    continue
                add_image(self.output_path + "\\" + pic_path, slide_pic)
                ppt_pic.save(self.output_path + "\\output_" + str(cnt) + "_" + str(curSlide) + ".pptx")
                curSlide += 1
            # 表格节点
            if node["type"] == 3:
                # 获取行列数
                table_data = node["content"]
                rows = len(table_data)
                cols = len(table_data[0])
                ppt_pic = Presentation(picture_path)
                slide_table = ppt_pic.slides[0]
                add_tables(slide_table, table_data, rows, cols)
                ppt_pic.save(self.output_path + "\\output_" + str(cnt) + "_" + str(curSlide) + ".pptx")
                curSlide += 1
        # 可能存在不足4个句子的节点 最后需要保存

        ppt.save(self.output_path + "\\output_" + str(cnt) + "_" + str(curSlide) + ".pptx")
        cnt = 0