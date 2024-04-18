from pptx import Presentation
import shutil
from Utils import Utils

utils = Utils()

if __name__=="__main__":
    template_path = utils.getTemplatePath() + "\\template01.pptx"
    output_path = utils.getRootPath() + "\output.pptx"
    shutil.copy(template_path, output_path)

    ppt = Presentation(output_path)
    # 封面页
    title_slide = ppt.slides[0]
    for shape in title_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                print(paragraph.text)
                if "互联网" in paragraph.text:
                    paragraph.text = "中文学术论文自动转幻灯片算法设计与系统实现"
    ppt.save(output_path)


    